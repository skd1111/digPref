"""文档解析：主流办公文档格式 → ParsedDocument（页 / 块 / 全局偏移）。

支持矩阵：
    pdf   pypdf 文本层；无文本层（扫描件）→ pypdfium2 栅格化 + RapidOCR 端侧识别回退
    docx  python-docx 段落 + 表格
    doc   Word 97-2003：Windows 下经 Word/WPS COM 转 docx；PK 魔数嗅探直解
    txt/md/csv  纯文本（UTF-8 / GBK）
    html  标准库 HTMLParser 提取正文文本
    xlsx  openpyxl 惰性导入（未安装 → 引导提示）
    pptx  python-pptx 惰性导入（未安装 → 引导提示）
    xls   旧二进制 Excel 不支持 → 引导另存为 xlsx
"""

from __future__ import annotations

import contextlib
import logging
import os
import re
import subprocess
import tempfile
import uuid
from html.parser import HTMLParser
from pathlib import Path
from typing import ClassVar

from agent.config import settings
from agent.doc_review.models import Block, DocFormat, Page, ParsedDocument, generate_id

logger = logging.getLogger(__name__)


class DocParseError(Exception):
    pass


_SUPPORTED_SUFFIXES = {
    ".pdf",
    ".docx",
    ".doc",
    ".txt",
    ".md",
    ".csv",
    ".html",
    ".htm",
    ".xlsx",
    ".pptx",
}

_FORMAT_BY_SUFFIX = {
    ".pdf": DocFormat.PDF,
    ".docx": DocFormat.DOCX,
    ".doc": DocFormat.DOC,
    ".txt": DocFormat.TXT,
    ".md": DocFormat.MD,
    ".csv": DocFormat.TXT,
    ".html": DocFormat.HTML,
    ".htm": DocFormat.HTML,
    ".xlsx": DocFormat.XLSX,
    ".pptx": DocFormat.PPTX,
}

_ZIP_MAGIC = b"PK\x03\x04"
_OLE2_MAGIC = b"\xd0\xcf\x11\xe0"


def _normalize_block_text(text: str) -> str:
    return re.sub(r"[ \t\r\f\v]+", " ", text).strip()


def _split_blocks(text: str) -> list[str]:
    raw = re.split(r"\n\s*\n", text)
    return [b for b in (_normalize_block_text(x) for x in raw) if b]


def _assemble(
    blocks_by_page: list[tuple[int, list[str]]],
    *,
    file_name: str,
    file_path: str,
    fmt: DocFormat,
) -> ParsedDocument:
    pages: list[Page] = []
    all_text: list[str] = []
    offset = 0
    for page_no, texts in blocks_by_page:
        blocks: list[Block] = []
        for i, text in enumerate(texts):
            block_id = f"p{page_no}b{i + 1}"
            start = offset
            end = start + len(text)
            blocks.append(Block(block_id=block_id, text=text, start=start, end=end))
            all_text.append(text)
            offset = end + 1  # 块之间以单个 '\n' 分隔
        pages.append(Page(page_no=page_no, blocks=blocks))
    return ParsedDocument(
        doc_id=generate_id(),
        file_name=file_name,
        file_path=file_path,
        format=fmt,
        page_count=len(pages),
        pages=pages,
        full_text="\n".join(all_text),
    )


def _parse_pdf(path: Path) -> list[tuple[int, list[str]]]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise DocParseError(f"PDF 解析失败: {exc}") from exc
    if not reader.pages:
        raise DocParseError("PDF 无页面")
    pages: list[tuple[int, list[str]]] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        blocks = _split_blocks(text)
        if blocks:
            pages.append((i, blocks))
    if pages:
        return pages
    # 无文本层 → 扫描件：按开关走 pypdfium2 栅格化 + RapidOCR 端侧识别回退（best-effort）
    return _ocr_scanned_pdf(path)


def _ocr_scanned_pdf(path: Path) -> list[tuple[int, list[str]]]:
    """扫描件 PDF 的 OCR 回退：仅当无文本层时触发（正常 PDF 不走）。

    纯本地端侧（pypdfium2 栅格化 + RapidOCR），数据不出域；开关关闭/依赖缺失/
    识别失败均静默返 []（上层 parse_document 报「未提取到文本」）。
    """
    if not settings.doc_review_pdf_ocr_enabled:
        return []
    try:
        from agent.image_processing.ocr import ocr_pdf_to_pages, rapidocr_available

        if not rapidocr_available():
            logger.info("scanned PDF OCR skipped: RapidOCR 不可用 file=%s", path.name)
            return []
        pages = ocr_pdf_to_pages(
            path,
            scale=float(settings.doc_review_pdf_ocr_scale),
            max_pages=int(settings.doc_review_pdf_ocr_max_pages),
        )
        if pages:
            logger.info("scanned PDF OCR fallback ok file=%s pages=%d", path.name, len(pages))
        return pages
    except Exception as exc:
        logger.warning("scanned PDF OCR fallback failed file=%s: %s", path.name, exc)
        return []


def _parse_docx(path: Path) -> list[tuple[int, list[str]]]:
    import docx
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = docx.Document(str(path))
    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, document).text or ""
            if text.strip():
                blocks.append(_normalize_block_text(text))
        elif child.tag == qn("w:tbl"):
            table = Table(child, document)
            for row in table.rows:
                for cell in row.cells:
                    cell_text = _normalize_block_text(cell.text or "")
                    if cell_text:
                        blocks.append(cell_text)
    return [(1, blocks)]


def _parse_text(path: Path) -> list[tuple[int, list[str]]]:
    raw = path.read_bytes()
    for encoding in ("utf-8", "gbk"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        raise DocParseError("文本编码不支持（支持 UTF-8 / GBK）")
    return [(1, _split_blocks(text))]


# ---- HTML（标准库，零依赖）----


class _HtmlTextExtractor(HTMLParser):
    _SKIP_TAGS: ClassVar[set[str]] = {"script", "style", "head", "noscript"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._chunks: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1
        if tag in ("p", "div", "br", "tr", "li", "h1", "h2", "h3", "h4", "h5", "h6", "table"):
            self._chunks.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0 and data.strip():
            self._chunks.append(data)

    def text(self) -> str:
        return "".join(self._chunks)


def _parse_html(path: Path) -> list[tuple[int, list[str]]]:
    raw = path.read_bytes()
    text: str | None = None
    for encoding in ("utf-8", "gbk"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        raise DocParseError("HTML 编码不支持（支持 UTF-8 / GBK）")
    extractor = _HtmlTextExtractor()
    try:
        extractor.feed(text)
    except Exception as exc:
        raise DocParseError(f"HTML 解析失败: {exc}") from exc
    blocks = _split_blocks(extractor.text())
    return [(1, blocks)]


# ---- xlsx / pptx（惰性导入，未装给引导提示）----


def _parse_xlsx(path: Path) -> list[tuple[int, list[str]]]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise DocParseError("Excel 解析需要 openpyxl：uv pip install openpyxl") from exc
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        raise DocParseError(f"Excel 解析失败: {exc}") from exc
    pages: list[tuple[int, list[str]]] = []
    for page_no, sheet in enumerate(wb.worksheets, start=1):
        blocks: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [_normalize_block_text("" if c is None else str(c)) for c in row]
            line = " | ".join(c for c in cells if c)
            if line:
                blocks.append(line)
        if blocks:
            pages.append((page_no, blocks))
    wb.close()
    return pages


def _parse_pptx(path: Path) -> list[tuple[int, list[str]]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocParseError("PPT 解析需要 python-pptx：uv pip install python-pptx") from exc
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise DocParseError(f"PPT 解析失败: {exc}") from exc
    pages: list[tuple[int, list[str]]] = []
    for page_no, slide in enumerate(prs.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = _normalize_block_text(shape.text_frame.text or "")
                if text:
                    blocks.append(text)
        if blocks:
            pages.append((page_no, blocks))
    return pages


# ---- .doc（Word 97-2003）：COM 转 docx ----

_WORD_CONVERT_PS = (
    "$ErrorActionPreference='Stop';"
    "$w = New-Object -ComObject Word.Application;"
    "$w.Visible = $false; $w.DisplayAlerts = 0;"
    "$d = $w.Documents.Open($env:EAIDE_DOC_SRC, $false, $true);"
    "$d.SaveAs([ref]$env:EAIDE_DOC_DST, [ref]16);"
    "$d.Close($false); $w.Quit()"
)


def _convert_doc_to_docx(path: Path) -> Path:
    """经 Word/WPS COM 把 .doc 转成 .docx（仅 Windows）。"""
    if os.name != "nt":
        raise DocParseError(
            ".doc 自动转换仅支持 Windows（本机请安装 Word 或 WPS），或另存为 .docx 后重新导入"
        )
    dst = Path(tempfile.gettempdir()) / f"eaide_docconv_{uuid.uuid4().hex}.docx"
    env = {**os.environ, "EAIDE_DOC_SRC": str(path.resolve()), "EAIDE_DOC_DST": str(dst)}
    try:
        proc = subprocess.run(
            [
                "powershell",
                "-NoProfile",
                "-ExecutionPolicy",
                "Bypass",
                "-Command",
                _WORD_CONVERT_PS,
            ],
            env=env,
            capture_output=True,
            text=True,
            timeout=120,
        )
    except FileNotFoundError as exc:
        raise DocParseError(".doc 转换失败：未找到 powershell") from exc
    except subprocess.TimeoutExpired as exc:
        raise DocParseError(
            ".doc 转换超时（120s）：Word/WPS 可能被弹窗阻塞，请另存为 .docx 后重新导入"
        ) from exc
    if proc.returncode != 0 or not dst.exists():
        detail = (proc.stderr or proc.stdout or "").strip().splitlines()[-1:] or ["未知错误"]
        raise DocParseError(
            f".doc 转换失败（需本机安装 Word 或 WPS）：{detail[0][:200]}；也可另存为 .docx 后重新导入"
        )
    return dst


def _parse_doc(path: Path) -> tuple[DocFormat, list[tuple[int, list[str]]]]:
    head = path.read_bytes()[:4]
    if head == _ZIP_MAGIC:
        # 实为 docx 被改名为 .doc —— 直接按 docx 解析
        return DocFormat.DOCX, _parse_docx(path)
    if head != _OLE2_MAGIC:
        raise DocParseError(".doc 文件头无法识别（非 OLE2 / ZIP），文件可能已损坏")
    converted = _convert_doc_to_docx(path)
    try:
        return DocFormat.DOC, _parse_docx(converted)
    finally:
        with contextlib.suppress(OSError):
            converted.unlink(missing_ok=True)


def parse_document(path: str | Path) -> ParsedDocument:
    p = Path(path)
    if not p.exists():
        raise DocParseError(f"文件不存在: {p}")
    suffix = p.suffix.lower()
    if suffix == ".xls":
        raise DocParseError(
            "不支持的格式: .xls（旧版 Excel 二进制格式）——请用 Excel 另存为 .xlsx 后重新导入"
        )
    if suffix not in _SUPPORTED_SUFFIXES:
        raise DocParseError(
            f"不支持的格式: {suffix}（支持 pdf/docx/doc/txt/md/csv/html/xlsx/pptx）"
        )
    fmt = _FORMAT_BY_SUFFIX[suffix]
    if fmt == DocFormat.PDF:
        by_page = _parse_pdf(p)
    elif fmt == DocFormat.DOCX:
        by_page = _parse_docx(p)
    elif fmt == DocFormat.DOC:
        fmt, by_page = _parse_doc(p)
    elif fmt == DocFormat.HTML:
        by_page = _parse_html(p)
    elif fmt == DocFormat.XLSX:
        by_page = _parse_xlsx(p)
    elif fmt == DocFormat.PPTX:
        by_page = _parse_pptx(p)
    else:
        by_page = _parse_text(p)
    if not by_page:
        raise DocParseError("未提取到文本（可能为扫描件且 OCR 未启用/不可用，或确实无文字内容）")
    return _assemble(by_page, file_name=p.name, file_path=str(p.resolve()), fmt=fmt)
